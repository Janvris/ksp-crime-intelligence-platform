from pptx import Presentation

prs = Presentation(r'D:\MY PROJECTS\Datathon\KSP_Datathon_2026_Submission.pptx')

def set_text(shape, new_text):
    if hasattr(shape, "text"):
        shape.text = new_text
        
# Slide 2: Problem statement
s2 = prs.slides[1]
for shape in s2.shapes:
    if shape.name == 'TextBox 60':
        set_text(shape, "Problem Statement:\nAdvanced Crime Analytics and Intelligence Platform for Karnataka State Police (KSP). Enable investigators, analysts, and policymakers to discover hidden criminal networks, predict emerging hotspots using AutoML, and understand socio-economic drivers of crime using interactive visualizations.")

# Slide 3: Brief about the solution
s3 = prs.slides[2]
for shape in s3.shapes:
    if shape.name == 'TextBox 60':
        set_text(shape, "A comprehensive crime analytics platform that provides actionable intelligence to KSP.\n\nCORE PILLARS:\n1. Network Analysis: Discover organized crime rings using Louvain community detection and Centrality algorithms.\n2. Predictive AI: District-level risk scoring and anomaly detection via Zia AutoML.\n3. Spatial Analysis: Real-time hotspot mapping using KDE.\n4. Socio-Economic Correlation: Correlate demographic data with crime rates to identify root causes.")

# Slide 5: Features
s5 = prs.slides[4]
for shape in s5.shapes:
    if shape.name == 'TextBox 72':
        set_text(shape, "MVP FEATURES\n\n1. Organized Crime Network Analysis\n   Graph of accused, FIRs, bank accounts, and phones. Algorithms: Louvain (communities), Centrality (key figures), and QuickML (MO similarity).\n\n2. Predictive AI & Anomaly Detection\n   District Risk Choropleth (Zia AutoML), real-time anomaly logs, and feature importance drivers.\n\n3. Spatial & Temporal Hotspots\n   Interactive Leaflet KDE heatmaps, combined with ECharts for temporal trend forecasting.")
    if shape.name == 'TextBox 73':
        set_text(shape, "CROSS-CUTTING\n\n4. Socio-Economic Analysis\n   Pearson correlation between crime and unemployment, literacy, etc.\n\n5. Explainable AI\n   Clear lineage trails for AI predictions, algorithms used, and parameters.\n\n6. Secure RBAC\n   Role-based dashboards for SCRB Analyst, District SP, and Investigator with varying data access.")

# Slide 7: UI Screens
s7 = prs.slides[6]
for shape in s7.shapes:
    if shape.name == 'TextBox 84':
        set_text(shape, "KEY UI MODULES IMPLEMENTED\n\nModule 1 — Main Dashboard\n   KPIs, Leaflet Spatial Heatmap, Temporal Trend Forecasts.\n\nModule 2 — Network Link Analysis\n   Seed-based entity search, force-directed graph with Accused/FIRs/Banks, and dynamic algorithm overlays (Louvain/Centrality).\n\nModule 3 — Predictive AI\n   Choropleth risk map, anomaly detection alerts, and feature importance charts.\n\nModule 4 — Socio-Economic Correlation\n   Scatter plots and correlation matrices matching demographics to crime types.")

# Slide 9: Technologies
s9 = prs.slides[8]
for shape in s9.shapes:
    if shape.name == 'TextBox 96':
        set_text(shape, "FRONTEND\n• HTML5 + Vanilla JavaScript (Lightweight & Fast)\n• Tailwind CSS (Responsive utility-first design)\n• ECharts (Complex graph viz, trends, scatter plots)\n• Leaflet + Leaflet.heat (Geospatial mapping)\n• Font Awesome (Icons)\n\nALGORITHMS IMPLEMENTED (Client-side)\n• Force-directed graph layouts\n• Pearson Correlation Coefficient\n• UI simulation of Louvain, Centrality, and QuickML MO matching.")
    if shape.name == 'TextBox 97':
        set_text(shape, "DATA\n• Static JSON Data Sets (Karnataka Crime Data)\n• Embedded socio-economic metrics\n\nINFRASTRUCTURE / HOSTING\n• Zoho Catalyst Slate (Web Client Hosting)\n• Future integration planned with Catalyst Data Store & QuickML for full backend operation.")

# Slide 10: Catalyst Services
s10 = prs.slides[9]
for shape in s10.shapes:
    if shape.name == 'TextBox 102':
        set_text(shape, "CATALYST SERVICES (Current & Planned)\n\n1. Slate — Web client hosting (Current)\n2. QuickML — For live MO matching (Planned)\n3. Zia AutoML — Crime forecasting models (Planned)\n4. Data Store — Relational crime DB (Planned)\n5. Authentication — Login & RBAC (Current simulation, Planned implementation)\n6. Functions — Agent logic & APIs (Planned)")
    if shape.name == 'TextBox 103':
        set_text(shape, "WHY CATALYST?\n\n• Rapid Deployment: Slate allows immediate deployment of the frontend.\n• Unified AI Ecosystem: Native access to Zia and QuickML makes transitioning from simulation to real AI seamless.\n• Serverless Scalability: Handles varying loads from state-wide police forces.")

# Slide 14: Links
s14 = prs.slides[13]
for shape in s14.shapes:
    if shape.name == 'TextBox 126':
        set_text(shape, "GitHub Public Repository:\nhttps://github.com/[your-team]/ksp-crime-analytics\n   (Contains: Frontend HTML/JS, Data JSONs, Catalyst configuration)\n\nDemo Video Link:\nhttps://[youtube/vimeo-link]\n   (Script: Problem -> Main Dashboard -> Network Analysis Algorithms -> Predictive/Socio-Economic AI)\n\nDeployed Link:\nhttps://ksp-crime.zoho.com (Catalyst Slate deployment)")

prs.save(r'D:\MY PROJECTS\Datathon\KSP_Datathon_2026_Submission_Updated.pptx')
print("Successfully updated presentation.")

import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation

prs = Presentation(r'D:\MY PROJECTS\Datathon\KSP_Datathon_2026_Submission.pptx')

for i, slide in enumerate(prs.slides):
    print(f"\n{'='*80}")
    print(f"SLIDE {i+1}")
    print(f"{'='*80}")
    for j, shape in enumerate(slide.shapes):
        if hasattr(shape, "text") and shape.text.strip():
            print(f"\n--- Shape '{shape.name}' ---")
            print(shape.text)

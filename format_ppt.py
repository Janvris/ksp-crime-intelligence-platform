from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation(r'D:\MY PROJECTS\Datathon\KSP_Datathon_2026_Submission.pptx')

def format_shape_text(shape, heading_text, body_text):
    if not hasattr(shape, "text_frame"):
        return
    tf = shape.text_frame
    tf.clear()
    
    # Heading
    if heading_text:
        p_head = tf.paragraphs[0]
        p_head.text = heading_text
        p_head.alignment = PP_ALIGN.LEFT
        for run in p_head.runs:
            run.font.name = 'Times New Roman'
            run.font.size = Pt(13)
            run.font.bold = True
    
    # Body
    if body_text:
        lines = body_text.split('\n')
        for i, line in enumerate(lines):
            if not heading_text and i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            # Justify longer lines, left align shorter ones
            p.alignment = PP_ALIGN.JUSTIFY if len(line) > 60 else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.name = 'Times New Roman'
                run.font.size = Pt(10)
                run.font.bold = False

# Map of slide index -> shape name -> (heading, body)
updates = {
    1: { # Slide 2: Problem
        'TextBox 60': ("Problem Statement", 
                       "Karnataka State Police faces significant challenges in analyzing siloed crime data, identifying organized criminal networks, and proactively predicting emerging hotspots. Existing systems rely on rigid dashboards and lack intuitive, AI-driven insights required for rapid decision-making.")
    },
    2: { # Slide 3: Solution
        'TextBox 60': ("KSP Crime Intelligence Platform", 
                       "An intelligent, centralized analytics hub that transforms raw crime data into actionable insights.\n\n"
                       "CORE PILLARS:\n"
                       "• Network Analysis: Uncover organized crime rings using Louvain community detection and Centrality algorithms.\n"
                       "• Predictive AI: Leverage Zia AutoML for district-level risk scoring and real-time anomaly detection.\n"
                       "• Spatial Intelligence: Map real-time hotspots using Kernel Density Estimation (KDE).\n"
                       "• Socio-Economic Insights: Correlate demographic data with crime rates to identify root causes.")
    },
    4: { # Slide 5: Features
        'TextBox 72': ("Key MVP Features", 
                       "1. Organized Crime Network Analysis\n"
                       "Interactive force-directed graphs linking Accused, FIRs, Bank Accounts, and Phones. Features UI-simulated Louvain (communities) and Centrality (key figures) algorithms.\n\n"
                       "2. Predictive AI & Anomaly Detection\n"
                       "Dynamic District Risk Choropleth maps, real-time anomaly alerts (e.g., vehicle theft surges), and feature importance drivers.\n\n"
                       "3. Spatial & Temporal Hotspots\n"
                       "Geospatial mapping with Leaflet KDE heatmaps, integrated with ECharts for analyzing temporal trend forecasts."),
        'TextBox 73': ("Cross-Cutting Capabilities", 
                       "4. Socio-Economic Correlation\n"
                       "Pearson correlation analysis linking crime types to demographics like unemployment and urbanization.\n\n"
                       "5. Explainable AI\n"
                       "Transparent lineage trails for all AI predictions, ensuring trust and accountability.\n\n"
                       "6. Secure Role-Based Access\n"
                       "Tailored dashboards and data access for SCRB Analysts, District SPs, and Investigators.")
    },
    6: { # Slide 7: UI Screens
        'TextBox 84': ("Key UI Modules Implemented", 
                       "Module 1: Main Dashboard\nHigh-level KPIs, spatial heatmaps, and temporal crime trend forecasts.\n\n"
                       "Module 2: Network Link Analysis\nSeed-based entity search powering dynamic graphs with algorithm overlays.\n\n"
                       "Module 3: Predictive AI\nChoropleth risk map, anomaly detection logs, and feature importance analysis.\n\n"
                       "Module 4: Socio-Economic Correlation\nInteractive scatter plots and correlation matrices matching demographics to crime rates.")
    },
    8: { # Slide 9: Technologies
        'TextBox 96': ("Frontend Technologies", 
                       "• HTML5 & Vanilla JavaScript: Lightweight, fast, and dependency-free.\n"
                       "• Tailwind CSS: Responsive, utility-first UI design for premium aesthetics.\n"
                       "• ECharts: Powerful rendering for complex graphs, trends, and scatter plots.\n"
                       "• Leaflet & Leaflet.heat: Interactive geospatial mapping.\n"
                       "• Font Awesome: Clean, professional iconography."),
        'TextBox 97': ("Algorithms & Infrastructure", 
                       "ALGORITHMS IMPLEMENTED\n"
                       "• Force-directed graph layouts for link analysis.\n"
                       "• Pearson Correlation Coefficient for socio-economic insights.\n"
                       "• Client-side simulations of Louvain, Centrality, and QuickML MO matching.\n\n"
                       "DATA & HOSTING\n"
                       "• Data: Static JSON Data Sets (Karnataka Crime Data) + embedded demographics.\n"
                       "• Hosting: Deployed on Zoho Catalyst Slate (Web Client Hosting).")
    },
    9: { # Slide 10: Catalyst Services
        'TextBox 102': ("Catalyst Services Utilization", 
                       "CURRENT IMPLEMENTATION\n"
                       "1. Slate: Web client hosting for the frontend application.\n\n"
                       "PLANNED INTEGRATION (Phase 2)\n"
                       "2. QuickML: Real-time MO similarity matching.\n"
                       "3. Zia AutoML: Time-series crime forecasting models.\n"
                       "4. Data Store: Robust relational crime database.\n"
                       "5. Authentication: Secure Login and RBAC implementation.\n"
                       "6. Functions: Serverless Agent logic and APIs."),
        'TextBox 103': ("Why Catalyst?", 
                       "• Rapid Deployment: Slate enabled immediate, hassle-free deployment of the frontend.\n"
                       "• Unified AI Ecosystem: Native access to Zia and QuickML ensures a seamless transition from prototype to production.\n"
                       "• Serverless Scalability: Designed to handle varying loads from state-wide police forces efficiently and cost-effectively.")
    },
    13: { # Slide 14: Links
        'TextBox 126': ("Project Links", 
                        "GitHub Public Repository:\n"
                        "https://github.com/[your-team]/ksp-crime-analytics\n"
                        "(Contains: Frontend HTML/JS, Data JSONs, Catalyst configuration)\n\n"
                        "Demo Video Link:\n"
                        "https://[youtube/vimeo-link]\n"
                        "(Walkthrough: Problem Statement -> Main Dashboard -> Network Algorithms -> Predictive AI)\n\n"
                        "Deployed Application:\n"
                        "https://ksp-crime.zoho.com (Catalyst Slate deployment)")
    }
}

for slide_idx, shape_updates in updates.items():
    slide = prs.slides[slide_idx]
    for shape in slide.shapes:
        if shape.name in shape_updates:
            heading, body = shape_updates[shape.name]
            format_shape_text(shape, heading, body)

# Apply global formatting to any remaining text shapes to ensure consistency
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame:
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.font.name = 'Times New Roman'
                    # If font size is None, we don't necessarily want to force it to 10 here 
                    # as it might ruin title slides, but we'll ensure the font family is consistent.

prs.save(r'D:\MY PROJECTS\Datathon\KSP_Datathon_2026_Submission_Formatted.pptx')
print("Successfully formatted and updated presentation.")
